# -*- coding: utf-8 -*-
"""
app/api/files.py - 文件上传和处理 API

提供文件上传、解析和处理功能
支持的文件类型：PDF、TXT、PPT
"""

from fastapi import APIRouter, UploadFile, File, HTTPException
from pydantic import BaseModel
from typing import List, Optional
import os
import uuid
import shutil
from openai import OpenAI

# 创建路由实例
router = APIRouter()

# 初始化 DeepSeek 客户端
deepseek_client = OpenAI(
    api_key=os.getenv("DEEPSEEK_API_KEY"),
    base_url="https://api.deepseek.com"
)

# 文件存储目录
UPLOAD_DIR = "uploads"
# 确保上传目录存在
os.makedirs(UPLOAD_DIR, exist_ok=True)


# ============================================
# 数据模型定义
# ============================================

class FileInfo(BaseModel):
    """文件信息模型"""
    id: str
    filename: str
    file_type: str
    size: int
    upload_time: str


class SummaryRequest(BaseModel):
    """文件总结请求"""
    file_id: str
    summary_type: str = "general"  # general, key_points, outline


class SummaryResponse(BaseModel):
    """文件总结响应"""
    summary: str
    key_points: List[str]
    file_id: str


# ============================================
# API 接口定义
# ============================================

@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """
    上传文件
    
    参数：
    - file: 上传的文件
    
    返回：
    - 文件信息和 ID
    """
    try:
        # 生成唯一 ID
        file_id = str(uuid.uuid4())
        
        # 获取文件扩展名
        filename = file.filename
        ext = os.path.splitext(filename)[1].lower()
        
        # 检查文件类型
        allowed_types = ['.pdf', '.txt', '.pptx', '.ppt']
        if ext not in allowed_types:
            raise HTTPException(
                status_code=400,
                detail=f"不支持的文件类型。仅支持: {', '.join(allowed_types)}"
            )
        
        # 生成新文件名（使用 UUID 避免重名）
        new_filename = f"{file_id}{ext}"
        file_path = os.path.join(UPLOAD_DIR, new_filename)
        
        # 保存文件
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        # 获取文件大小
        file_size = os.path.getsize(file_path)
        
        return {
            "id": file_id,
            "filename": filename,
            "file_type": ext[1:],  # 去掉点
            "size": file_size,
            "message": "文件上传成功"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"上传失败: {str(e)}")


@router.get("/list")
async def list_files():
    """
    获取已上传文件列表
    
    返回：
    - 文件列表
    """
    # TODO: 从数据库读取文件列表
    return {"files": []}


@router.get("/{file_id}")
async def get_file(file_id: str):
    """
    获取指定文件的信息
    
    参数：
    - file_id: 文件 ID
    
    返回：
    - 文件信息
    """
    # 查找文件
    for ext in ['.pdf', '.txt', '.pptx', '.ppt']:
        file_path = os.path.join(UPLOAD_DIR, f"{file_id}{ext}")
        if os.path.exists(file_path):
            return {
                "id": file_id,
                "filename": f"file{ext}",
                "size": os.path.getsize(file_path),
                "path": file_path
            }
    
    raise HTTPException(status_code=404, detail="文件不存在")


@router.delete("/{file_id}")
async def delete_file(file_id: str):
    """
    删除指定文件
    
    参数：
    - file_id: 文件 ID
    
    返回：
    - 成功消息
    """
    deleted = False
    # 尝试删除各种扩展名的文件
    for ext in ['.pdf', '.txt', '.pptx', '.ppt']:
        file_path = os.path.join(UPLOAD_DIR, f"{file_id}{ext}")
        if os.path.exists(file_path):
            os.remove(file_path)
            deleted = True
            break
    
    if deleted:
        return {"message": "文件删除成功"}
    else:
        raise HTTPException(status_code=404, detail="文件不存在")


@router.post("/summarize", response_model=SummaryResponse)
async def summarize_file(request: SummaryRequest):
    """
    AI 总结文件内容
    
    参数：
    - request: 总结请求，包含文件 ID 和总结类型
    
    返回：
    - 总结结果
    """
    # 查找文件
    file_path = None
    file_ext = None
    for ext in ['.pdf', '.txt', '.pptx', '.ppt']:
        temp_path = os.path.join(UPLOAD_DIR, f"{request.file_id}{ext}")
        if os.path.exists(temp_path):
            file_path = temp_path
            file_ext = ext
            break
    
    if not file_path:
        raise HTTPException(status_code=404, detail="文件不存在")
    
    # 提取文件内容
    content = ""
    try:
        if file_ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()[:8000]  # 限制长度
        elif file_ext == '.pdf':
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages[:10]:  # 最多10页
                    text = page.extract_text()
                    if text:
                        content += text + "\n"
                content = content[:8000]
        elif file_ext in ['.pptx', '.ppt']:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides[:20]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            content = content[:8000]
    except Exception as e:
        return SummaryResponse(
            summary=f"读取文件失败: {str(e)}",
            key_points=["文件读取错误"],
            file_id=request.file_id
        )
    
    if not content.strip():
        return SummaryResponse(
            summary="文件内容为空或无法提取",
            key_points=["无法提取内容"],
            file_id=request.file_id
        )
    
    # 调用 DeepSeek API 进行总结
    try:
        prompt = f"""请帮我总结以下文档内容，提取关键要点。

文档内容：
{content}

请用中文回复，格式如下：
1. 简要总结（100字以内）
2. 关键要点（3-5个）
"""

        response = deepseek_client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个专业的学习助手，擅长总结文档内容和提取关键知识点。"},
                {"role": "user", "content": prompt}
            ],
            temperature=0.5,
            max_tokens=1000
        )
        
        summary_text = response.choices[0].message.content
        
        # 简单解析关键点
        key_points = []
        lines = summary_text.split('\n')
        for line in lines:
            if line.strip() and (line.strip()[0].isdigit() or '•' in line or '-' in line):
                key_points.append(line.strip()[:100])
        
        return SummaryResponse(
            summary=summary_text,
            key_points=key_points[:5] if key_points else ["详见上方总结"],
            file_id=request.file_id
        )
        
    except Exception as e:
        return SummaryResponse(
            summary=f"总结失败: {str(e)}",
            key_points=["AI 服务调用失败"],
            file_id=request.file_id
        )


@router.get("/{file_id}/content")
async def get_file_content(file_id: str):
    """
    获取文件内容（文本提取）
    
    参数：
    - file_id: 文件 ID
    
    返回：
    - 文件文本内容
    """
    # 查找文件
    file_path = None
    file_ext = None
    for ext in ['.pdf', '.txt', '.pptx', '.ppt']:
        temp_path = os.path.join(UPLOAD_DIR, f"{file_id}{ext}")
        if os.path.exists(temp_path):
            file_path = temp_path
            file_ext = ext
            break
    
    if not file_path:
        raise HTTPException(status_code=404, detail="文件不存在")
    
    try:
        # 根据文件类型提取内容
        if file_ext == '.txt':
            # TXT 文件直接读取
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        elif file_ext == '.pdf':
            # PDF 文件使用 pdfplumber 提取
            import pdfplumber
            content = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        content += text + "\n"
        elif file_ext in ['.pptx', '.ppt']:
            # PPT 文件使用 python-pptx 提取
            from pptx import Presentation
            content = ""
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        else:
            content = "不支持的文件类型"
        
        return {
            "content": content,
            "file_id": file_id,
            "file_type": file_ext[1:]
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"读取文件失败: {str(e)}")
